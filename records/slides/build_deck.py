#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""恩送り AI実践セミナー  PowerPoint 生成（slides_data.py を参照）"""
import os
import qrcode
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import slides_data as D

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
os.makedirs(ASSETS, exist_ok=True)

def C(t):  # tuple -> RGBColor
    return RGBColor(t[0], t[1], t[2])

NAVY, NAVY2, GOLD, GOLD_LT = C(D.NAVY), C(D.NAVY2), C(D.GOLD), C(D.GOLD_LT)
PAPER, INK, MUTED, WHITE, CARD, CARD_LN, BODY_LT = (
    C(D.PAPER), C(D.INK), C(D.MUTED), C(D.WHITE), C(D.CARD), C(D.CARD_LN), C(D.BODY_LT))
SERIF, SANS = "Noto Serif CJK JP", "Noto Sans CJK JP"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def make_qr(data, path):
    qr = qrcode.QRCode(box_size=12, border=2,
                       error_correction=qrcode.constants.ERROR_CORRECT_M)
    qr.add_data(data); qr.make(fit=True)
    qr.make_image(fill_color=(0x14, 0x21, 0x3D), back_color="white").save(path)
    return path


def solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background(); shape.shadow.inherit = False


def grad(shape, c1, c2, angle=90):
    shape.fill.gradient()
    st = shape.fill.gradient_stops
    st[0].position, st[0].color.rgb = 0.0, c1
    st[1].position, st[1].color.rgb = 1.0, c2
    try: shape.fill.gradient_angle = angle
    except Exception: pass
    shape.line.fill.background(); shape.shadow.inherit = False


def send_back(slide, shape):
    sp = shape._element; sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def bg(slide, color, gradient=False):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    if gradient: grad(r, NAVY2, NAVY)
    else: solid(r, color)
    send_back(slide, r); return r


def rect(slide, x, y, w, h, color, line=None, lw=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    solid(s, color)
    if line is not None:
        s.line.color.rgb = line; s.line.width = lw or Pt(1)
    return s


def ea(run, font):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", font)


def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, ls=1.0, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        p.space_before = Pt(0); p.line_spacing = ls
        for (t, font, size, color, bold) in para:
            r = p.add_run(); r.text = t
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = bold
            ea(r, font)
    return tb


def page_no(slide, n):
    text(slide, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4),
         [[("%02d / 24" % n, SANS, 10, MUTED, False)]], align=PP_ALIGN.RIGHT)


def kicker(slide, num, label):
    text(slide, Inches(0.9), Inches(0.62), Inches(9), Inches(0.5),
         [[("%02d" % num, SANS, 14, GOLD, True), ("   " + label, SANS, 13, MUTED, False)]])


def heading(slide, title, size=34):
    text(slide, Inches(0.9), Inches(1.05), Inches(11.6), Inches(1.0),
         [[(title, SANS, size, INK, True)]])
    rect(slide, Inches(0.93), Inches(1.95), Inches(1.1), Inches(0.05), GOLD)


# ---------- レンダラ ----------
def render_cover(s, d):
    bg(s, None, gradient=True)
    rect(s, Inches(0.9), Inches(2.1), Inches(0.06), Inches(3.3), GOLD)
    text(s, Inches(1.2), Inches(1.55), Inches(11), Inches(0.6), [[(d["org"], SANS, 16, GOLD_LT, False)]])
    text(s, Inches(1.15), Inches(2.2), Inches(11), Inches(2.2), [[(d["title"], SERIF, 88, WHITE, True)]])
    text(s, Inches(1.2), Inches(4.5), Inches(11.5), Inches(1.0), [[(d["lead"], SANS, 26, GOLD, True)]])
    text(s, Inches(1.2), Inches(5.45), Inches(11.5), Inches(0.8), [[(d["sub"], SANS, 18, BODY_LT, False)]])
    text(s, Inches(1.2), Inches(6.7), Inches(11.5), Inches(0.5), [[(d["presenter"], SANS, 13, RGBColor(0x9A,0xA4,0xBC), False)]])


def render_statement(s, d):
    bg(s, None, gradient=True)
    rect(s, Inches(0.95), Inches(2.5), Inches(0.9), Inches(0.06), GOLD)
    text(s, Inches(0.95), Inches(1.7), Inches(11), Inches(0.6), [[(d["kicker"], SANS, 16, GOLD_LT, False)]])
    text(s, Inches(0.95), Inches(2.9), Inches(11.6), Inches(2.4),
         [[(ln, SERIF, 46, WHITE, True)] for ln in d["big"]], ls=1.15, space_after=4)
    yy = 2.9 + 0.95 * len(d["big"]) + 0.35
    if d.get("sub"):
        text(s, Inches(0.95), Inches(yy), Inches(11.6), Inches(2.0),
             [[(ln, SANS, 18, BODY_LT, False)] for ln in d["sub"]], ls=1.3, space_after=4)
    if d.get("footnote"):
        text(s, Inches(0.95), Inches(6.65), Inches(11.6), Inches(0.6),
             [[(d["footnote"], SANS, 13.5, GOLD, True)]])
    page_no(s, d["num"])


def render_content(s, d):
    bg(s, PAPER); kicker(s, d["num"], d["label"]); heading(s, d["title"])
    top = 2.35
    if d.get("body"):
        runs = []
        for t, lv in d["body"]:
            if lv == 0:
                runs.append([("", SANS, 6, INK, False)])
                runs.append([(t, SANS, 19, NAVY, True)])
            else:
                runs.append([("▸  ", SANS, 16, GOLD, True), (t, SANS, 17, INK, False)])
        text(s, Inches(0.95), Inches(top), Inches(11.6), Inches(4.4), runs, ls=1.12, space_after=7)
    if d.get("two_col"):
        colw = Inches(5.55)
        ctop = 2.3
        for ci, (head, items) in enumerate(d["two_col"]):
            x = Inches(0.95) if ci == 0 else Inches(6.85)
            card = rect(s, x, Inches(ctop), colw, Inches(3.9), CARD, line=CARD_LN, lw=Pt(1))
            text(s, x + Inches(0.3), Inches(ctop + 0.22), colw - Inches(0.6), Inches(0.6),
                 [[(head, SANS, 18, NAVY, True)]])
            rect(s, x + Inches(0.32), Inches(ctop + 0.72), Inches(0.7), Inches(0.04), GOLD)
            sz = 14.5 if len(items) <= 7 else 13
            gap = 5 if len(items) <= 7 else 2
            runs = [[("・ ", SANS, sz - 0.5, GOLD, True), (it, SANS, sz, INK, False)] for it in items]
            text(s, x + Inches(0.32), Inches(ctop + 0.92), colw - Inches(0.64), Inches(2.8),
                 runs, ls=1.1, space_after=gap)
    if d.get("note"):
        rect(s, Inches(0.95), Inches(6.35), Inches(11.45), Inches(0.72), NAVY)
        text(s, Inches(1.25), Inches(6.35), Inches(10.9), Inches(0.72),
             [[(d["note"], SANS, 14.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    page_no(s, d["num"])


def render_work(s, d, qr_path):
    bg(s, PAPER); kicker(s, d["num"], d["label"]); heading(s, d["title"], size=32)
    text(s, Inches(0.95), Inches(2.45), Inches(7.7), Inches(3.5),
         [[(ln, SANS, 17.5, INK, False)] for ln in d["lines"]], ls=1.35, space_after=10)
    qx, qy, qs = Inches(9.4), Inches(2.5), Inches(2.85)
    card = rect(s, qx - Inches(0.28), qy - Inches(0.28), qs + Inches(0.56), qs + Inches(1.15), CARD,
                line=GOLD, lw=Pt(1.5))
    s.shapes.add_picture(qr_path, qx, qy, qs, qs)
    text(s, qx - Inches(0.28), qy + qs + Inches(0.02), qs + Inches(0.56), Inches(0.8),
         [[(d["qr_caption"], SANS, 14, NAVY, True)], [("QRから回答", SANS, 11.5, MUTED, False)]],
         align=PP_ALIGN.CENTER, space_after=2)
    text(s, Inches(0.95), Inches(6.55), Inches(7.7), Inches(0.5),
         [[("※ QRは仮URL。実フォーム作成後に差し替えます。", SANS, 11.5, MUTED, False)]])
    page_no(s, d["num"])


def render_prompt(s, d):
    bg(s, PAPER); kicker(s, d["num"], d["label"])
    text(s, Inches(0.9), Inches(0.98), Inches(11.6), Inches(0.8), [[(d["title"], SANS, 30, INK, True)]])
    rect(s, Inches(0.93), Inches(1.78), Inches(1.1), Inches(0.05), GOLD)
    rect(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(4.85), NAVY)
    text(s, Inches(1.25), Inches(2.3), Inches(10.9), Inches(0.6), [[(d["intro"], SANS, 15, GOLD_LT, True)]])
    text(s, Inches(1.25), Inches(2.95), Inches(10.9), Inches(3.8),
         [[(b, SANS, 13.5, RGBColor(0xE8,0xEC,0xF5), False)] for b in d["blocks"]], ls=1.2, space_after=3)
    page_no(s, d["num"])


# ---------- 生成 ----------
qr_paths = {
    "kickoff": make_qr(D.URL_KICKOFF, os.path.join(ASSETS, "qr_kickoff.png")),
    "survey":  make_qr(D.URL_SURVEY,  os.path.join(ASSETS, "qr_survey.png")),
}

for d in D.SLIDES:
    s = prs.slides.add_slide(BLANK)
    k = d["kind"]
    if k == "cover":       render_cover(s, d)
    elif k == "statement": render_statement(s, d)
    elif k == "content":   render_content(s, d)
    elif k == "work":      render_work(s, d, qr_paths[d["qr"]])
    elif k == "prompt":    render_prompt(s, d)

out = os.path.join(HERE, "恩送り_AI実践セミナー.pptx")
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
